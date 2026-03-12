#!/usr/bin/env python3
"""SecureHub v2 プレゼンテーション PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)
PRIMARY_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
PRIMARY_DARK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_LIGHT = RGBColor(0x22, 0xD3, 0xEE)
DANGER = RGBColor(0xEF, 0x44, 0x44)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x3A, 0x5F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=14, default_color=TEXT_LIGHT):
    """lines: list of (text, font_size, color, bold, alignment) or str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold, align = line, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else False
            align = line[4] if len(line) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_card(slide, left, top, width, height, title, desc, title_color=WHITE, icon_text="", border_color=None):
    shape = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=border_color or RGBColor(0x3B, 0x82, 0xF6))
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 f"{icon_text} {title}" if icon_text else title, font_size=14, color=title_color, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                 desc, font_size=11, color=TEXT_MUTED)


def add_stat_card(slide, left, top, width, height, number, label, source="", number_color=DANGER):
    shape = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=RGBColor(0x3B, 0x82, 0xF6))
    add_text_box(slide, left, top + Inches(0.2), width, Inches(0.5),
                 number, font_size=36, color=number_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.8), width - Inches(0.3), Inches(0.4),
                 label, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    if source:
        add_text_box(slide, left + Inches(0.15), top + Inches(1.15), width - Inches(0.3), Inches(0.3),
                     source, font_size=8, color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)


# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Shield icon (text-based)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.9),
             "\U0001F6E1", font_size=60, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(1.0),
             "SecureHub", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.1), SLIDE_W, Inches(0.6),
             "セキュリティ文化を変える", font_size=24, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.8), SLIDE_W, Inches(0.5),
             "インシデント報告 & セキュリティ文化醸成プラットフォーム", font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(0.4),
             "2026年3月", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "皆さん、こんにちは。本日は「セキュアハブ」をご紹介します。組織のセキュリティ文化そのものを変革するプラットフォームです。"


# ===== SLIDE 2: Current Problems =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "⚠ セキュリティの現状", font_size=32, color=ACCENT_LIGHT, bold=True)
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(10), Inches(0.4),
             "組織が直面している深刻な課題", font_size=14, color=TEXT_MUTED)

# Stat cards (2x2)
card_w = Inches(5.8)
card_h = Inches(1.5)
x1, x2 = Inches(0.6), Inches(6.8)
y1, y2 = Inches(1.7), Inches(3.5)

add_stat_card(slide, x1, y1, card_w, card_h, "70%+", "リスクある行動を自認する従業員", "Proofpoint State of the Phish 2024", DANGER)
add_stat_card(slide, x2, y1, card_w, card_h, "88%", "盗まれた認証情報が関与するデータ侵害", "Verizon DBIR 2025", DANGER)
add_stat_card(slide, x1, y2, card_w, card_h, "$4.4M", "データ侵害の世界平均コスト", "IBM Cost of a Data Breach 2025", WARNING)
add_stat_card(slide, x2, y2, card_w, card_h, "年1回", "多くの企業のセキュリティ教育の実態", "形式的なeラーニング受講のみ", PRIMARY_LIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = "現状を見てみましょう。従業員の70%以上がリスク行動を自覚しながら利便性を優先。データ侵害の88%は認証情報の窃取が原因で、平均コストは440万ドルです。しかし多くの企業の教育は年1回のeラーニングだけ。これでは行動は変わりません。"


# ===== SLIDE 3: Gap Analysis =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "🔍 Gap分析：根本課題", font_size=32, color=ACCENT_LIGHT, bold=True)
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(10), Inches(0.4),
             "現状とあるべき姿の間に横たわる4つの課題", font_size=14, color=TEXT_MUTED)

gaps = [
    ("報告忌避文化", "ミスの報告が叱責につながる恐怖。インシデントが隠蔽され、検知が遅延する"),
    ("経営層の関与不足", "セキュリティをIT部門の問題と矮小化。予算・人員の優先度が低い"),
    ("形式的な教育", "年1回のeラーニング。受講率だけがKPI、行動変容を測定しない"),
    ("当事者意識の欠如", "「セキュリティはIT部門の仕事」。全員の責任という意識が根付いていない"),
]

for i, (title, desc) in enumerate(gaps):
    col = i % 2
    row = i // 2
    left = Inches(0.6) + col * Inches(6.2)
    top = Inches(1.6) + row * Inches(1.6)
    # Red left border card
    shape = add_shape(slide, left, top, Inches(5.8), Inches(1.3), fill_color=BG_CARD)
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(1.3))
    border.fill.solid()
    border.fill.fore_color.rgb = DANGER
    border.line.fill.background()
    add_text_box(slide, left + Inches(0.25), top + Inches(0.1), Inches(5.3), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.5), Inches(5.3), Inches(0.7),
                 desc, font_size=11, color=TEXT_MUTED)

# Highlight box
highlight = add_shape(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.7),
                      fill_color=RGBColor(0x0A, 0x2E, 0x4E), border_color=ACCENT)
add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5),
             "→ 共通するのは「組織文化」の問題 — 技術だけでは解決できない",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "Gap分析から4つの根本課題が見えました。報告忌避文化、経営層の関与不足、形式的な教育、当事者意識の欠如。すべてが技術ではなく「組織文化」の問題です。文化を変える仕組みが必要です。"


# ===== SLIDE 4: Comparison with Seculio (NEW) =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "⚖ 既存サービスとの違い", font_size=32, color=ACCENT_LIGHT, bold=True)
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(10), Inches(0.4),
             "教育型サービスとの差別化 — SecureHubが埋める空白領域", font_size=14, color=TEXT_MUTED)

# Table header
table_left = Inches(0.6)
table_top = Inches(1.6)
col_widths = [Inches(3.0), Inches(4.8), Inches(4.8)]
row_height = Inches(0.5)
header_height = Inches(0.55)

headers = ["機能・観点", "既存サービス（セキュリオ等）", "SecureHub"]
header_colors = [RGBColor(0x1A, 0x30, 0x6B), RGBColor(0x33, 0x33, 0x44), RGBColor(0x0A, 0x3D, 0x5C)]

rows_data = [
    ("主軸", "教育・訓練・テスト（知識定着）", "報告文化・脅威共有・当事者意識（文化変革）", True),
    ("教育・テスト", "✓ 100種以上の教材、リテラシースコア", "― 対象外（既存サービスと併用想定）", False),
    ("フィッシング訓練", "✓ 標的型攻撃メール訓練", "― 対象外", False),
    ("インシデント報告", "フィッシング報告機能", "✓ ワンクリック+匿名+AI自動トリアージ+感謝FB", True),
    ("脅威情報のAI要約配信", "― なし", "✓ IPA/JPCERT等からAI自動要約＋全社配信", True),
    ("匿名Q&A", "― なし", "✓ AIチャットボット＋ナレッジベース化", True),
    ("チャンピオン制度", "― なし", "✓ 部門別リーダー管理・表彰・コミュニティ", True),
]

# Draw headers
x = table_left
for i, (header, hcolor) in enumerate(zip(headers, header_colors)):
    shape = add_shape(slide, x, table_top, col_widths[i], header_height, fill_color=hcolor)
    add_text_box(slide, x + Inches(0.1), table_top + Inches(0.05), col_widths[i] - Inches(0.2), header_height,
                 header, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    x += col_widths[i]

# Draw rows
for r, (label, seculio, securehub, highlight) in enumerate(rows_data):
    y = table_top + header_height + r * row_height
    bg = RGBColor(0x12, 0x1E, 0x35) if r % 2 == 0 else RGBColor(0x0F, 0x17, 0x2A)
    x = table_left
    for c, (text, w) in enumerate(zip([label, seculio, securehub], col_widths)):
        shape = add_shape(slide, x, y, w, row_height, fill_color=bg)
        tc = WHITE if c == 0 else (ACCENT_LIGHT if (c == 2 and highlight) else TEXT_LIGHT)
        fb = True if c == 0 else (True if (c == 2 and highlight) else False)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), row_height - Inches(0.1),
                     text, font_size=10, color=tc, bold=fb)
        x += w

# Highlight box
highlight_box = add_shape(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.65),
                          fill_color=RGBColor(0x0A, 0x2E, 0x4E), border_color=ACCENT)
add_text_box(slide, Inches(1.5), Inches(5.88), Inches(10.3), Inches(0.5),
             "🧩 「教育で知識を付ける」既存サービスと、「文化を変える」SecureHubは補完関係",
             font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "ここで、既存サービスとの違いを整理します。セキュリオなどの既存サービスは「教育・テスト・訓練」で知識を定着させるのが主軸です。一方、セキュアハブは「報告文化・脅威共有・当事者意識」という文化変革が主軸。脅威情報のAI要約配信、匿名Q&A、チャンピオン制度は既存サービスにない独自機能です。つまり、両者は競合ではなく補完関係。教育と文化変革の両輪で、組織を強くします。"


# ===== SLIDE 5: Concept =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "💡 SecureHubのコンセプト", font_size=32, color=ACCENT_LIGHT, bold=True)
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(10), Inches(0.4),
             "「セキュリティは全員の責任」を日常の行動に変える", font_size=14, color=TEXT_MUTED)

pillars = [
    ("🚩 報告文化", "ワンクリックで報告。匿名対応。\n報告者を称賛し、\n心理的安全性を確保する", DANGER),
    ("📰 脅威共有", "最新の脅威情報を全員に\nリアルタイム配信。\nAI要約で誰でも理解できる", WARNING),
    ("👥 当事者意識", "チャンピオン制度で各部門に\nリーダーを。セキュリティを\n「みんなのこと」に", SUCCESS),
]

for i, (title, desc, accent_color) in enumerate(pillars):
    left = Inches(0.6) + i * Inches(4.2)
    top = Inches(1.7)
    card = add_shape(slide, left, top, Inches(3.8), Inches(2.8), fill_color=BG_CARD, border_color=accent_color)
    add_text_box(slide, left, top + Inches(0.3), Inches(3.8), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.0), Inches(3.2), Inches(1.5),
                 desc, font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Highlight box
highlight_box = add_shape(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.7),
                          fill_color=RGBColor(0x0A, 0x2E, 0x4E), border_color=ACCENT)
add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.5),
             "❝ セキュリティを「やらされること」から「みんなで取り組むこと」へ ❞",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "SecureHubは3本の柱で構成されています。報告文化、脅威共有、当事者意識。目指すのは、セキュリティを「やらされること」から「みんなで取り組むこと」に変えることです。"


# ===== SLIDE 6: 5 Features =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "📦 5つの主要機能", font_size=32, color=ACCENT_LIGHT, bold=True)

features = [
    ("🚩", "ワンクリック報告", "メールアドイン連携\n匿名報告対応\nAI自動トリアージ"),
    ("📡", "脅威情報フィード", "自動収集・AI要約\n週次ダイジェスト\n緊急アラート"),
    ("💬", "匿名Q&A", "匿名質問投稿\nAIチャットボット\nナレッジベース化"),
    ("🏆", "チャンピオン管理", "活動ログ管理\nバッジ・表彰\nコミュニティ"),
    ("📊", "経営ダッシュボード", "リスクスコア\n部門別比較\n自動レポート"),
]

for i, (icon, title, desc) in enumerate(features):
    left = Inches(0.4) + i * Inches(2.5)
    top = Inches(1.3)
    card = add_shape(slide, left, top, Inches(2.3), Inches(2.5), fill_color=BG_CARD, border_color=RGBColor(0x3B, 0x82, 0xF6))
    add_text_box(slide, left, top + Inches(0.15), Inches(2.3), Inches(0.5),
                 icon, font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.7), Inches(2.3), Inches(0.35),
                 title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.2), Inches(2.0), Inches(1.2),
                 desc, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Flow
flow_steps = ["👁 気づく", "→", "🚩 報告する", "→", "🧠 学ぶ", "→", "🔗 共有する", "→", "🛡 守る"]
flow_text = "    ".join(flow_steps)
flow_bg = add_shape(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(0.8), fill_color=BG_CARD, border_color=RGBColor(0x3B, 0x82, 0xF6))
add_text_box(slide, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.6),
             flow_text, font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "主要機能は5つ。ワンクリック報告、脅威情報フィード、匿名Q&A、チャンピオン管理、経営ダッシュボード。「気づく、報告する、学ぶ、共有する、守る」。このサイクルが日常になることで、セキュリティが文化になります。"


# ===== SLIDE 7: Feature Detail =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "🚩 報告機能 & 📡 脅威フィード", font_size=32, color=ACCENT_LIGHT, bold=True)

# Report card
card1 = add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.8), fill_color=BG_CARD, border_color=RGBColor(0x3B, 0x82, 0xF6))
add_text_box(slide, Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.4),
             "🚩 ワンクリックインシデント報告", font_size=16, color=WHITE, bold=True)
report_items = [
    "● メールアドイン: Outlook/Gmailにワンクリック報告ボタン",
    "● 匿名報告: 実名/匿名を選択可能",
    "● AI自動トリアージ: 緊急度を自動分類",
    "● 感謝フィードバック: 報告者に感謝 + 対応状況通知",
]
add_multiline_text(slide, Inches(0.9), Inches(1.9), Inches(5.2), Inches(2.0), report_items, default_size=11, default_color=TEXT_MUTED)

# Threat feed card
card2 = add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), fill_color=BG_CARD, border_color=RGBColor(0x3B, 0x82, 0xF6))
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.4),
             "📡 セキュリティニュース・脅威情報", font_size=16, color=WHITE, bold=True)
threat_items = [
    "● 自動収集: IPA, JPCERT/CC等からRSS/API連携",
    "● AI要約: 非技術者向けに平易に要約",
    "● 週次ダイジェスト: Slack/Teams/メールで配信",
    "● 緊急アラート: 重大脅威のプッシュ通知",
]
add_multiline_text(slide, Inches(7.1), Inches(1.9), Inches(5.2), Inches(2.0), threat_items, default_size=11, default_color=TEXT_MUTED)

# Flow
flow_labels = ["📧 不審メール発見", "→", "☝ ワンクリック報告", "→", "🤖 AI自動トリアージ", "→", "❤ 感謝フィードバック"]
flow_bg = add_shape(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.7), fill_color=BG_CARD, border_color=ACCENT)
add_text_box(slide, Inches(0.6), Inches(4.55), Inches(12.1), Inches(0.6),
             "    ".join(flow_labels), font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "報告機能の核心は「ハードルを下げる」こと。ワンクリックで匿名報告、AIが自動トリアージ。報告後は感謝メッセージが届き、「報告してよかった」と思える体験を設計しています。脅威情報フィードでは、IPAやJPCERT等の情報をAIが自動要約。専門知識がなくても、最新の脅威が分かりやすく届きます。"


# ===== SLIDE 8: Demo - Report Screen =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "🖥 デモ：インシデント報告画面", font_size=32, color=ACCENT_LIGHT, bold=True)

# Browser mock
browser = add_shape(slide, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.5), fill_color=RGBColor(0x1A, 0x1A, 0x2E))
# Browser bar
bar = add_shape(slide, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.5), fill_color=RGBColor(0x16, 0x21, 0x3E))
add_text_box(slide, Inches(1.3), Inches(1.25), Inches(10), Inches(0.4),
             "● ● ●  🔒 securehub.example.com/report", font_size=10, color=TEXT_MUTED)

# Form content header
add_text_box(slide, Inches(1.4), Inches(1.9), Inches(8), Inches(0.4),
             "🚩 インシデント報告", font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(3.6), Inches(1.95), Inches(8), Inches(0.3),
             "安心して報告してください。あなたの報告が組織を守ります。", font_size=10, color=TEXT_MUTED)

# Category section
add_text_box(slide, Inches(1.4), Inches(2.5), Inches(4), Inches(0.3),
             "カテゴリを選択", font_size=12, color=TEXT_LIGHT, bold=True)

categories = ["📧 不審メール", "📞 不審な電話", "🚪 物理セキュリティ", "💾 データ関連", "💻 デバイス紛失", "... その他"]
for i, cat in enumerate(categories):
    col = i % 3
    row = i // 3
    left = Inches(1.4) + col * Inches(1.8)
    top = Inches(2.9) + row * Inches(0.65)
    is_selected = (i == 0)
    bc = ACCENT if is_selected else RGBColor(0x3B, 0x82, 0xF6)
    fc = RGBColor(0x0A, 0x2E, 0x4E) if is_selected else RGBColor(0x0F, 0x17, 0x2A)
    cat_shape = add_shape(slide, left, top, Inches(1.6), Inches(0.55), fill_color=fc, border_color=bc)
    tc = ACCENT_LIGHT if is_selected else TEXT_MUTED
    add_text_box(slide, left, top + Inches(0.1), Inches(1.6), Inches(0.35),
                 cat, font_size=10, color=tc, alignment=PP_ALIGN.CENTER)

# Anonymous toggle
add_text_box(slide, Inches(1.4), Inches(4.4), Inches(4), Inches(0.3),
             "🔘 匿名で報告する", font_size=11, color=ACCENT_LIGHT)

# Detail section
add_text_box(slide, Inches(7.0), Inches(2.5), Inches(4), Inches(0.3),
             "詳細（任意）", font_size=12, color=TEXT_LIGHT, bold=True)
detail_box = add_shape(slide, Inches(7.0), Inches(2.9), Inches(4.8), Inches(1.8), fill_color=RGBColor(0x0F, 0x17, 0x2A), border_color=RGBColor(0x3B, 0x82, 0xF6))
add_text_box(slide, Inches(7.2), Inches(3.0), Inches(4.4), Inches(1.5),
             "件名に「アカウント確認のお願い」とある\nメールが届きました。\nリンク先のURLが不審です...", font_size=10, color=TEXT_MUTED)

# Attachment
attach_box = add_shape(slide, Inches(7.0), Inches(4.9), Inches(4.8), Inches(0.45), fill_color=RGBColor(0x0F, 0x17, 0x2A), border_color=RGBColor(0x3B, 0x82, 0xF6))
add_text_box(slide, Inches(7.2), Inches(4.95), Inches(4.4), Inches(0.35),
             "📎 スクリーンショットを添付（任意）", font_size=10, color=TEXT_MUTED)

# Submit button
submit = add_shape(slide, Inches(7.0), Inches(5.6), Inches(2.2), Inches(0.5), fill_color=PRIMARY, border_color=ACCENT)
add_text_box(slide, Inches(7.0), Inches(5.65), Inches(2.2), Inches(0.4),
             "✈ 報告を送信", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "実際の報告画面です。カテゴリを選んで、匿名トグルをオンにすれば完全匿名。詳細は任意入力、スクリーンショットも添付可能。送信ボタンを押すだけで完了です。とにかくシンプルに。心理的ハードルを下げることにこだわりました。"


# ===== SLIDE 9: Demo - Dashboard =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "📊 デモ：経営層ダッシュボード", font_size=32, color=ACCENT_LIGHT, bold=True)

# Browser mock
browser = add_shape(slide, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.5), fill_color=RGBColor(0x1A, 0x1A, 0x2E))
bar = add_shape(slide, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.5), fill_color=RGBColor(0x16, 0x21, 0x3E))
add_text_box(slide, Inches(1.3), Inches(1.25), Inches(10), Inches(0.4),
             "● ● ●  🔒 securehub.example.com/dashboard", font_size=10, color=TEXT_MUTED)

add_text_box(slide, Inches(1.4), Inches(1.9), Inches(6), Inches(0.4),
             "📊 セキュリティ文化ダッシュボード", font_size=16, color=WHITE, bold=True)
add_text_box(slide, Inches(8.5), Inches(1.95), Inches(3), Inches(0.3),
             "2026年3月 | 📥 PDF出力", font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

# KPI cards
kpis = [
    ("72", "ヒューマンリスクスコア", "↑ +8 先月比", SUCCESS),
    ("156", "月間報告件数", "↑ +23% 先月比", ACCENT),
    ("47min", "平均報告時間", "↓ -15min 先月比", WARNING),
    ("68%", "エンゲージメント率", "↑ +5% 先月比", PRIMARY_LIGHT),
]

for i, (value, label, trend, color) in enumerate(kpis):
    left = Inches(1.4) + i * Inches(2.7)
    top = Inches(2.5)
    card = add_shape(slide, left, top, Inches(2.4), Inches(1.3), fill_color=RGBColor(0x0F, 0x17, 0x2A), border_color=RGBColor(0x3B, 0x82, 0xF6))
    add_text_box(slide, left, top + Inches(0.1), Inches(2.4), Inches(0.5),
                 value, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.65), Inches(2.4), Inches(0.25),
                 label, font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.9), Inches(2.4), Inches(0.25),
                 trend, font_size=9, color=SUCCESS, alignment=PP_ALIGN.CENTER)

# Chart area
chart_bg = add_shape(slide, Inches(1.4), Inches(4.1), Inches(10.5), Inches(2.3), fill_color=RGBColor(0x0F, 0x17, 0x2A), border_color=RGBColor(0x3B, 0x82, 0xF6))
add_text_box(slide, Inches(1.6), Inches(4.2), Inches(8), Inches(0.3),
             "📊 部門別報告件数（月次推移）", font_size=10, color=TEXT_MUTED)

depts = ["営業", "開発", "管理", "人事", "企画", "製造", "IT", "法務", "経理", "総務"]
heights = [0.35, 0.6, 0.8, 0.45, 0.7, 0.55, 0.9, 0.4, 0.5, 0.65]
bar_width = Inches(0.85)
for i, (dept, h) in enumerate(zip(depts, heights)):
    left = Inches(1.7) + i * Inches(1.0)
    bar_h = Inches(h * 1.8)
    bar_top = Inches(6.1) - bar_h
    bar_shape = add_shape(slide, left, bar_top, bar_width, bar_h, fill_color=ACCENT)
    bar_shape.line.fill.background()
    add_text_box(slide, left, Inches(6.15), bar_width, Inches(0.2),
                 dept, font_size=8, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "経営層向けダッシュボードです。リスクスコア72点、月間報告156件、平均報告時間47分。すべての数値が改善トレンド。部門別の報告状況も一目で分かります。このダッシュボードで経営層のコミットメントが高まります。"


# ===== SLIDE 10: Closing =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
             "🚀 期待される効果とロードマップ", font_size=32, color=ACCENT_LIGHT, bold=True)

# Effect cards
effects = [
    ("🚩", "0 → 50+件/月", "インシデント報告件数"),
    ("⏱", "数日 → 1時間以内", "平均報告時間"),
    ("📈", "Lv.2 → Lv.4", "SANS成熟度レベル（2年目標）"),
]

for i, (icon, value, label) in enumerate(effects):
    left = Inches(0.6) + i * Inches(4.2)
    card = add_shape(slide, left, Inches(1.3), Inches(3.8), Inches(1.3), fill_color=BG_CARD, border_color=ACCENT)
    add_text_box(slide, left, Inches(1.35), Inches(3.8), Inches(0.4),
                 icon, font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(1.75), Inches(3.8), Inches(0.4),
                 value, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.15), Inches(3.8), Inches(0.3),
                 label, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Roadmap
phases = [
    ("Phase 1", "0〜2ヶ月", "MVP開発\nワンクリック報告\n基本ダッシュボード", PRIMARY_LIGHT),
    ("Phase 2", "2〜4ヶ月", "ニュースフィード\nQ&A機能\nSlack/Teams連携", ACCENT),
    ("Phase 3", "4〜5.5ヶ月", "チャンピオン管理\n表彰・バッジ\n経営ダッシュボード", SUCCESS),
    ("Phase 4", "5.5〜7ヶ月", "AIチャットボット\n自動トリアージ\nニュースAI要約", WARNING),
]

for i, (name, period, items, color) in enumerate(phases):
    left = Inches(0.4) + i * Inches(3.2)
    top = Inches(3.0)
    card = add_shape(slide, left, top, Inches(2.9), Inches(2.3), fill_color=BG_CARD)
    # Top border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.9), Inches(0.05))
    border.fill.solid()
    border.fill.fore_color.rgb = color
    border.line.fill.background()
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), Inches(2.6), Inches(0.3),
                 name, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.45), Inches(2.6), Inches(0.25),
                 period, font_size=10, color=color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.8), Inches(2.6), Inches(1.4),
                 items, font_size=10, color=TEXT_MUTED)

# Closing message
add_text_box(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.5),
             "🛡 SecureHub — セキュリティを「みんなのこと」に",
             font_size=20, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = "期待される効果。インシデント報告がゼロから月50件以上に。報告時間は数日から1時間以内に。SANS成熟度をLevel 4に引き上げます。開発は4フェーズ、7ヶ月。MVPは2ヶ月で稼働。SecureHubで、セキュリティを「みんなのこと」に。組織の文化を、一緒に変えていきましょう。ご清聴ありがとうございました。"


# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "demo_presentation_v2.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
